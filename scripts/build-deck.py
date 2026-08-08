#!/usr/bin/env python3
"""
Build the RAAHI Loom/demo deck: docs/RAAHI-loom-deck.pptx

    python3 scripts/build-deck.py

Deliberately a script rather than a hand-made file, for the same reason the
ElevenLabs agent is provisioned by script: the deck has to be regenerable after
a number changes. Test counts, tool counts and latency figures all moved several
times during the build, and a deck that silently disagrees with the README is
worse than no deck.

Writes to RAAHI-loom-deck.pptx and never touches RAAHI-demo-deck.pptx, which is
hand-authored by the team.

Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Emirates-adjacent palette, matching src/routes/talk.js
RED = RGBColor(0xD7, 0x19, 0x21)
RED_DARK = RGBColor(0xA3, 0x12, 0x1A)
GOLD = RGBColor(0xC3, 0x9C, 0x57)
INK = RGBColor(0x1C, 0x1C, 0x1C)
SOFT = RGBColor(0x5C, 0x60, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CANVAS = RGBColor(0xF7, 0xF7, 0xF8)
FONT = 'Helvetica Neue'

LOGO = 'docs/assets/raahi-symbol.png'
HERO = 'docs/assets/raahi-hero.jpg'
QR = 'docs/assets/raahi-qr.png'
ARCH = 'docs/assets/architecture-stack.png'
SEQ = 'docs/assets/demo-call-sequence.png'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(bg=CANVAS):
    s = prs.slides.add_slide(BLANK)
    bgshape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bgshape.fill.solid()
    bgshape.fill.fore_color.rgb = bg
    bgshape.line.fill.background()
    return s


def full_bleed(s, img=HERO):
    """Cover the slide with artwork, cropping to 16:9 rather than distorting."""
    try:
        s.shapes.add_picture(img, 0, 0, width=SW, height=SH)
        return True
    except Exception:
        return False


def scrim(s, x, y, w, h, color=RGBColor(0x12, 0x14, 0x18), alpha=None):
    """
    Solid dark panel so text stays legible over photography.

    Was translucent, via an a:alpha element injected straight into the fill XML
    because python-pptx exposes no alpha API. The package validated and
    reopened cleanly, but that is a private-schema edit and renderers flagged
    the file. A solid panel is one property, reads as deliberate design, and
    cannot be mis-parsed anywhere. `alpha` is accepted and ignored so callers
    need not change.
    """
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    return box


def band(s, y=0, h=Inches(0.10), color=RED):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SW, h)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    return r


def text(s, txt, x, y, w, h, size=18, color=INK, bold=False,
         align=PP_ALIGN.LEFT, space_after=6, line=1.25, caps=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = txt.split('\n') if isinstance(txt, str) else txt
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line
        run = p.add_run()
        run.text = ln.upper() if caps else ln
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return tb


def header(s, kicker, title):
    band(s)
    try:
        s.shapes.add_picture(LOGO, Inches(0.55), Inches(0.42), height=Inches(0.42))
    except Exception:
        pass
    text(s, kicker, Inches(1.15), Inches(0.44), Inches(9), Inches(0.3),
         size=11, color=RED, bold=True, caps=True)
    text(s, title, Inches(0.55), Inches(0.95), Inches(12.2), Inches(0.9),
         size=32, color=INK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.72), Inches(0.62), Pt(3))
    r.fill.solid(); r.fill.fore_color.rgb = GOLD
    r.line.fill.background()


def bullets(s, items, x=Inches(0.65), y=Inches(2.05), w=Inches(12.0), size=17, gap=0.62):
    for i, (head, body) in enumerate(items):
        yy = y + Inches(i * gap)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x, yy + Inches(0.09), Inches(0.11), Inches(0.11))
        dot.fill.solid(); dot.fill.fore_color.rgb = GOLD
        dot.line.fill.background()
        tb = s.shapes.add_textbox(x + Inches(0.28), yy - Inches(0.04), w, Inches(0.55))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.line_spacing = 1.2
        r1 = p.add_run(); r1.text = head + '  '
        r1.font.size = Pt(size); r1.font.bold = True
        r1.font.color.rgb = INK; r1.font.name = FONT
        if body:
            r2 = p.add_run(); r2.text = body
            r2.font.size = Pt(size - 1); r2.font.color.rgb = SOFT; r2.font.name = FONT


# Speaker notes are collected here and written to a companion markdown file
# rather than embedded as PowerPoint notes slides.
#
# python-pptx's notes_slide generation produces a notesMaster that macOS's
# renderer and PowerPoint both reject — verified by bisection: a one-slide deck
# with a single line of notes fails to render, the identical deck without notes
# is fine. That is what made the deck look corrupted. A companion file is also
# more useful while recording, since it can sit on a second screen instead of
# behind the slide.
SPEAKER_NOTES = []


def notes(s, txt):
    SPEAKER_NOTES.append(txt)


# ─────────────────────────────────────────────────────────── 1 · Title
s = slide(WHITE)
has_hero = full_bleed(s)
if has_hero:
    # Left-hand scrim: the artwork's subject sits right of centre, so the copy
    # goes left where the frame is quietest.
    scrim(s, 0, 0, Inches(6.9), SH, alpha=72000)
try:
    s.shapes.add_picture(LOGO, Inches(0.85), Inches(1.45), height=Inches(0.95))
except Exception:
    pass
text(s, 'RAAHI', Inches(0.85), Inches(2.6), Inches(6.0), Inches(1.1),
     size=64, color=WHITE if has_hero else RED, bold=True)
text(s, 'Y O U R   W A Y   F O R W A R D', Inches(0.9), Inches(3.72), Inches(6.0), Inches(0.4),
     size=13, color=GOLD)
text(s, 'A voice copilot that reads the page airlines\nactually publish — and tells a stranded\npassenger what to do next.',
     Inches(0.85), Inches(4.3), Inches(5.9), Inches(1.4),
     size=19, color=WHITE if has_hero else INK, line=1.4)
text(s, 'Siraj  ·  Astha  ·  Farman\nBUiD Voice Agents Hackathon, Dubai  ·  8 August 2026',
     Inches(0.85), Inches(6.05), Inches(6.0), Inches(0.8),
     size=12, color=RGBColor(0xC9, 0xCC, 0xD2) if has_hero else SOFT, line=1.35)
band(s, y=SH - Inches(0.09), h=Inches(0.09))
notes(s, "Hold 3 seconds on the artwork. Don't read the slide. Go straight to the problem.")

# ─────────────────────────────────────────────────────────── 2 · Problem
s = slide()
header(s, 'The problem', 'Two sentences decide whether she should get in a taxi')
text(s, '“The UAE will not allow entry to travellers who have recently been in the\nDemocratic Republic of Congo, Uganda, or South Sudan, unless the traveller\nhas been outside of these countries for more than 21 days.”',
     Inches(0.9), Inches(2.15), Inches(11.5), Inches(1.5), size=20, color=RED_DARK, line=1.4)
text(s, '“The entry and transit restrictions apply to all travellers, even those\narriving by indirect routings.”',
     Inches(0.9), Inches(3.7), Inches(11.5), Inches(1.0), size=20, color=RED_DARK, line=1.4)
bullets(s, [
    ('Not in any database.', 'Prose on emirates.com/help/travel-updates — changed in June, will change again.'),
    ('Buried.', 'Between a lounge refurbishment notice and Schengen biometrics.'),
    ('Decisive.', 'Kampala → Dubai → London is blocked. The queue to ask is measured in hours.'),
], y=Inches(4.95), size=16, gap=0.55)
notes(s, "Screen: the live advisory page, scrolled to these sentences. Let them be read. "
         "Say: 'these are not in any database — they're prose on a web page.'")

# ─────────────────────────────────────────────────────────── 3 · Demo
s = slide()
header(s, 'The demo', 'One booking reference → five live checks → one decision')
text(s, '“My booking reference is K seven X two M nine — should I leave for the airport?”',
     Inches(0.65), Inches(2.1), Inches(12.0), Inches(0.5), size=21, color=RED, bold=True)
bullets(s, [
    ('transit_rules', 'live advisory — restriction keys off ORIGIN, not destination'),
    ('disruption_status', 'live advisory — is London itself affected'),
    ('entry_requirements', 'live advisory + gov.uk + europa.eu — UK ETA'),
    ('track_aircraft', 'live ADS-B transponder — where the aircraft physically is'),
    ('weather_ops', 'live METAR — Dubai station conditions'),
], y=Inches(2.85), size=16, gap=0.52)
box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(5.6), Inches(12.0), Inches(1.2))
box.fill.solid(); box.fill.fore_color.rgb = WHITE
box.line.color.rgb = GOLD
text(s, 'Answered in 1.98 seconds, fanned out in parallel.        Control case: P3L8QK — Mumbai → London.\nSame destination. Opposite answer. Because the restriction depends on where you have BEEN.',
     Inches(0.9), Inches(5.78), Inches(11.6), Inches(1.0), size=16, color=INK, line=1.35)
notes(s, "Screen: phone on /talk. Say the PNR out loud. Then P3L8QK. "
         "The contrast IS the pitch — same destination, opposite answer.")

# ─────────────────────────────────────────────────────────── 4 · Architecture
s = slide()
header(s, 'Architecture', 'Five layers, and failure flows down gracefully')
try:
    s.shapes.add_picture(ARCH, Inches(0.5), Inches(2.0), width=Inches(12.3))
except Exception:
    text(s, '[docs/assets/architecture-stack.png]', Inches(0.65), Inches(2.4), Inches(12), Inches(0.5), size=14, color=SOFT)
notes(s, "Point at the ONE arrow pointing back in — the context.dev monitor webhook. "
         "Everything else is request/response; that is a push.")

# ─────────────────────────────────────────────────────────── 5 · Call flow
s = slide()
header(s, 'The call flow', 'PNR in, decision out — with the parallel fan-out')
try:
    s.shapes.add_picture(SEQ, Inches(1.9), Inches(1.95), height=Inches(5.25))
except Exception:
    text(s, '[docs/assets/demo-call-sequence.png]', Inches(0.65), Inches(2.4), Inches(12), Inches(0.5), size=14, color=SOFT)
notes(s, "Zoom on the parallel block while saying: one booking reference, five live sources, one decision.")

# ─────────────────────────────────────────────────────────── 6 · Intelligence
s = slide()
header(s, 'The intelligence', 'What a keyword match would get wrong')
bullets(s, [
    ('Two-sided matching', 'Restrictions key off where you have BEEN. Destination-only checking tells the Kampala passenger she is fine. She is not.'),
    ('Sentence windows', 'The qualifier that changes the answer sits in a neighbouring sentence, so every match reads one either side.'),
    ('Effective-from ≠ until', '“effective 6 June 2026 (until further notice)” contains a date that is not an end date. Returns open_ended, not a lifted restriction.'),
    ('Main-clause classification', '“If you do NOT need a visa… you WILL need an ETA.” The negation is in the subordinate clause. Filing that as an exemption gets someone denied boarding.'),
    ('Alias resolution', 'Callers say “DRC”, “the Congo”, “Kampala”. The page says “Democratic Republic of Congo”, “Uganda”.'),
], y=Inches(2.1), size=15, gap=0.88)
notes(s, "This is the slide that separates us from a chatbot with a search box. "
         "Every one of these is pinned by a test against verbatim page text.")

# ─────────────────────────────────────────────────────────── 7 · Sponsor tech
s = slide()
header(s, 'Sponsor technology', 'context.dev and ElevenLabs, used to their edges')
text(s, 'context.dev — 4 capabilities', Inches(0.65), Inches(2.05), Inches(6.0), Inches(0.4),
     size=15, color=RED, bold=True, caps=True)
bullets(s, [
    ('scrape/markdown', '4 verified sources'),
    ('web/search', 'allowlisted to 3 official domains'),
    ('monitors', 'semantic diff every 15 min, signed webhook'),
    ('scrape/html', 'evaluated; markdown sufficed'),
], x=Inches(0.65), y=Inches(2.55), w=Inches(5.6), size=14, gap=0.5)
text(s, 'ElevenLabs — both integration paths', Inches(6.9), Inches(2.05), Inches(6.0), Inches(0.4),
     size=15, color=RED, bold=True, caps=True)
bullets(s, [
    ('12 webhook tools', 'schedule & policy'),
    ('Native MCP server', 'live ADS-B tracking'),
    ('ASR keyword boost', 'callsigns, codes, Arabic terms'),
    ('Language presets', 'Arabic + Hindi, switches mid-call'),
    ('Guardrails', 'focus, injection, content, 2 custom'),
    ('Eval criteria', 'every call scored post-hoc'),
], x=Inches(6.9), y=Inches(2.55), w=Inches(5.9), size=14, gap=0.5)
text(s, 'Plus: adsb.lol ADS-B (keyless, ODbL) · aviationweather.gov METAR · Render · Node 18 / Express',
     Inches(0.65), Inches(6.25), Inches(12.0), Inches(0.5), size=14, color=SOFT)
notes(s, "Pick ONE to demo live: change detection, airspace snapshot, the card-number guardrail, or Arabic. Not all four.")

# ─────────────────────────────────────────────────────────── 8 · Honesty
s = slide()
header(s, 'What is actually live', 'Every response carries its own source field')
bullets(s, [
    ('Live on every call', 'disruption_status · transit_rules · entry_requirements · weather_ops · recent_changes · travel_intel · track_aircraft · airspace_snapshot'),
    ('Split-source, labelled', 'flight_status — ADS-B position live, schedule static. Two source fields so the live half cannot lend credibility to the static half.'),
    ('Static by design', 'policy_lookup — entitlements do not change hourly.'),
], y=Inches(2.1), size=15, gap=0.95)
text(s, 'What we do NOT claim', Inches(0.65), Inches(5.0), Inches(12), Inches(0.4),
     size=15, color=RED, bold=True, caps=True)
bullets(s, [
    ('No real PNR lookup.', 'emirates.com/manage-booking is auth-gated — we scraped it and got a login redirect. The booking store is a labelled stub; everything downstream is live.'),
    ('No live gates or delays.', 'Commercial data. ADS-B gives position, not schedule, and we label it as position.'),
    ('No push into a live turn.', 'The monitor pushes to us in minutes; the agent learns on its next tool call.'),
], y=Inches(5.45), size=14, gap=0.52)
notes(s, "Say this before a judge asks. Overclaiming is explicitly penalised; naming limits reads as engineering maturity.")

# ─────────────────────────────────────────────────────────── 9 · Reliability
s = slide()
header(s, 'Reliability & guardrails', 'The demo has to survive venue wifi')
bullets(s, [
    ('Never a 500', 'Any throw becomes 200 + a spoken-friendly message. A failed tool call is a voice agent going silent mid-sentence.'),
    ('live → cache → static → apology', 'Four-step degradation on every source call, with the honest label attached.'),
    ('58 tests, key blank', 'The offline path is the one tested hardest.'),
    ('Two guardrail layers', 'Platform guardrails outside the LLM + prompt rules. Verified adversarially: 5/5.'),
    ('Backend redaction', 'Cards, CVVs, IBANs, passports scrubbed from logs — Luhn-checked so ticket numbers survive.'),
    ('Injection defanged', 'Scraped prose is untrusted input handed to an LLM. Neutralised at the choke point.'),
], y=Inches(2.1), size=15, gap=0.75)
notes(s, "Reserve slide. Use if judges probe on robustness or security.")

# ─────────────────────────────────────────────────────────── 10 · Bugs found
s = slide()
header(s, 'Two bugs worth admitting', 'Both found by testing, both pinned')
bullets(s, [
    ('The agent invented an all-clear.', 'Given an unreadable tool result it said “clear to travel, proceed to the airport.” The backend said BLOCKED. It would have sent her to be refused at the gate. Now: an unreadable result is never good news, and it will not soften under a repeated question.'),
    ('One sentence could flip blocked → allowed.', 'The origin-side check asked a document-level question of a ±1 sentence window. An intervening sentence pushed the proof out and reversed the answer — reachable by Emirates simply editing their page. Scope is now read document-wide.'),
], y=Inches(2.2), size=16, gap=2.1)
notes(s, "Optional but strong. Shows we tested adversarially rather than demoed happy paths.")

# ─────────────────────────────────────────────── 11 · Vision (photographic)
s = slide(WHITE)
has_hero = full_bleed(s)
if has_hero:
    scrim(s, 0, Inches(4.42), SW, Inches(3.08), alpha=78000)
text(s, 'One question. Every layer. Two seconds.',
     Inches(0.85), Inches(4.72), Inches(11.6), Inches(0.7),
     size=34, color=WHITE if has_hero else INK, bold=True)
text(s, 'Voice · live advisory prose · entry rules · real aircraft · change detection — resolved into one\nsentence a passenger can act on, in the language they are frightened in.',
     Inches(0.85), Inches(5.5), Inches(11.6), Inches(1.1),
     size=17, color=RGBColor(0xDC, 0xDF, 0xE4) if has_hero else SOFT, line=1.4)
band(s, y=SH - Inches(0.09), h=Inches(0.09))
notes(s, "Use as the emotional beat before the call to action, or as the closing hold.")

# ─────────────────────────────────────────────────────────── 12 · Try it
s = slide(WHITE)
band(s, h=Inches(0.14))
band(s, y=Inches(0.14), h=Pt(3), color=GOLD)
text(s, 'Scan and talk to Raahi', Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.7),
     size=40, color=INK, align=PP_ALIGN.CENTER)
text(s, 'No app. A phone browser and a microphone.', Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.4),
     size=17, color=SOFT, align=PP_ALIGN.CENTER)
try:
    s.shapes.add_picture(QR, Inches(5.42), Inches(2.15), height=Inches(3.5))
except Exception:
    pass
text(s, 'irops-copilot-backend.onrender.com/talk', Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.4),
     size=17, color=RED_DARK, align=PP_ALIGN.CENTER, bold=True)
text(s, 'Try: “My booking reference is K seven X two M nine.”     github.com/sirajtechy/Airlines-Voice-Agent',
     Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4), size=13, color=SOFT, align=PP_ALIGN.CENTER)
notes(s, "End here. Leave it on screen while you close.")

prs.save('docs/RAAHI-loom-deck.pptx')

TITLES = [
    'Title — hero artwork',
    'The problem — two sentences',
    'The demo — one PNR, five live checks',
    'Architecture — five layers',
    'The call flow — parallel fan-out',
    'The intelligence — what keywords miss',
    'Sponsor technology — context.dev x ElevenLabs',
    'What is actually live',
    'Reliability and guardrails  (reserve)',
    'Two bugs worth admitting  (reserve)',
    'Vision — closing beat',
    'Try it — scan the QR',
]

with open('docs/RAAHI-deck-notes.md', 'w') as fh:
    fh.write('# Speaker notes — RAAHI demo deck\n\n')
    fh.write('Companion to [RAAHI-loom-deck.pptx](RAAHI-loom-deck.pptx). Generated by\n')
    fh.write('`scripts/build-deck.py` — edit the script, not this file.\n\n')
    fh.write('Notes live here rather than inside the .pptx because python-pptx\'s notes-slide\n')
    fh.write('output is rejected by PowerPoint and by macOS\'s renderer, which made the deck\n')
    fh.write('appear corrupted. Keep this open on a second screen while recording.\n\n')
    for i, (title, note) in enumerate(zip(TITLES, SPEAKER_NOTES), 1):
        fh.write(f'### Slide {i} · {title}\n\n{note}\n\n')

print(f'docs/RAAHI-loom-deck.pptx — {len(prs.slides._sldIdLst)} slides, no embedded notes')
print(f'docs/RAAHI-deck-notes.md  — {len(SPEAKER_NOTES)} speaker notes')
